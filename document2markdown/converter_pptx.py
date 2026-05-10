"""PPTX converter — maps python-pptx presentation elements to IR blocks.

Supported mappings:

* Slide title placeholder  → :class:`HeadingBlock` (level 2)
* Body / content text      → :class:`ParagraphBlock`
* Raster images            → :class:`EmbeddedAsset` + :class:`ImageBlock`
* EMF / WMF / EPS shapes   → :class:`VectorConverter` → :class:`EmbeddedAsset` + :class:`ImageBlock`
* Unrenderable elements    → :class:`UnsupportedBlock`

Slides are processed in presentation order.  Within each slide, the title
placeholder (if present) is emitted first as an H2 heading, followed by all
body text runs, then any extracted images.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from document2markdown.config import RASTER_DPI
from document2markdown.converter_base import BaseConverter
from document2markdown.converter_vector import VectorConverter, VectorConversionError
from document2markdown.document_model import (
    ConversionResult,
    EmbeddedAsset,
    HeadingBlock,
    ImageBlock,
    ParagraphBlock,
    UnsupportedBlock,
)

if TYPE_CHECKING:
    from document2markdown.document_model import Block

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# MIME / extension helpers (mirrors converter_docx.py)
# ---------------------------------------------------------------------------

_CONTENT_TYPE_TO_EXT: dict[str, str] = {
    "image/png":     ".png",
    "image/jpeg":    ".jpg",
    "image/jpg":     ".jpg",
    "image/gif":     ".gif",
    "image/bmp":     ".bmp",
    "image/tiff":    ".tif",
    "image/webp":    ".webp",
    "image/svg+xml": ".svg",
    "image/x-emf":   ".emf",
    "image/x-wmf":   ".wmf",
    "image/emf":     ".emf",
    "image/wmf":     ".wmf",
    "image/x-eps":   ".eps",
    "application/postscript": ".eps",
}

_VECTOR_CONTENT_TYPES: dict[str, str] = {
    "image/x-emf":            "emf",
    "image/emf":              "emf",
    "image/x-wmf":            "wmf",
    "image/wmf":              "wmf",
    "image/x-eps":            "eps",
    "application/postscript": "eps",
}

# python-pptx placeholder type constants
_TITLE_PLACEHOLDER_TYPES = {1, 13, 15}  # TITLE, CENTER_TITLE, VERTICAL_TITLE


def _content_type_to_ext(content_type: str) -> str:
    ct = content_type.split(";")[0].strip().lower()
    return _CONTENT_TYPE_TO_EXT.get(ct, ".bin")


def _content_type_to_vector_format(content_type: str) -> str | None:
    ct = content_type.split(";")[0].strip().lower()
    return _VECTOR_CONTENT_TYPES.get(ct)


# ---------------------------------------------------------------------------
# Main converter
# ---------------------------------------------------------------------------


class PPTXConverter(BaseConverter):
    """Convert a ``.pptx`` file to a :class:`ConversionResult`.

    Parameters
    ----------
    raster_dpi:
        DPI used when rasterizing vector graphics as a last resort.
    """

    def __init__(self, raster_dpi: int = RASTER_DPI) -> None:
        self._vector = VectorConverter(raster_dpi=raster_dpi)

    # ------------------------------------------------------------------
    # Public interface
    # ------------------------------------------------------------------

    def convert(self, source_path: Path) -> ConversionResult:
        """Convert *source_path* (.pptx) to a :class:`ConversionResult`."""
        prs = Presentation(str(source_path))

        blocks: list[Block] = []
        embedded: list[EmbeddedAsset] = []
        warnings: list[str] = []

        for slide_index, slide in enumerate(prs.slides):
            self._convert_slide(
                slide, slide_index, blocks, embedded, warnings
            )

        return ConversionResult(blocks=blocks, embedded=embedded, warnings=warnings)

    # ------------------------------------------------------------------
    # Slide conversion
    # ------------------------------------------------------------------

    def _convert_slide(
        self,
        slide,
        slide_index: int,
        blocks: list[Block],
        embedded: list[EmbeddedAsset],
        warnings: list[str],
    ) -> None:
        """Process a single slide and append blocks to *blocks*."""
        title_text: str | None = None

        # --- Extract title first ---
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                ph = shape.placeholder_format
            except ValueError:
                ph = None
            if ph is not None and ph.type in _TITLE_PLACEHOLDER_TYPES:
                title_text = shape.text_frame.text.strip()
                break

        if title_text:
            blocks.append(HeadingBlock(level=2, text=title_text))

        # --- Process all shapes in order ---
        for shape in slide.shapes:
            # Skip the title placeholder (already handled above)
            try:
                ph = shape.placeholder_format
            except ValueError:
                ph = None
            if ph is not None and ph.type in _TITLE_PLACEHOLDER_TYPES:
                continue

            self._convert_shape(shape, slide_index, blocks, embedded, warnings)

    # ------------------------------------------------------------------
    # Shape conversion
    # ------------------------------------------------------------------

    def _convert_shape(
        self,
        shape,
        slide_index: int,
        blocks: list[Block],
        embedded: list[EmbeddedAsset],
        warnings: list[str],
    ) -> None:
        """Convert a single shape and append resulting blocks."""
        shape_type = shape.shape_type

        # ---- Text frame (body text, content placeholders, text boxes) ----
        if shape.has_text_frame:
            self._extract_text_frame(shape, blocks)
            return

        # ---- Picture / raster image ----
        if shape_type == MSO_SHAPE_TYPE.PICTURE:
            block = self._extract_picture(shape, embedded, warnings)
            if block:
                blocks.append(block)
            return

        # ---- Group shape — recurse into children ----
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._convert_shape(child, slide_index, blocks, embedded, warnings)
            return

        # ---- OLE object (may contain EMF/WMF) ----
        if shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
            block = self._extract_ole_object(shape, embedded, warnings)
            if block:
                blocks.append(block)
            else:
                desc = f"slide {slide_index + 1}: unrenderable OLE object '{shape.name}'"
                warnings.append(desc)
                blocks.append(UnsupportedBlock(description=desc))
            return

        # ---- Auto-shape / freeform with no text — skip silently ----
        if shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.LINE,
        ):
            return

        # ---- Anything else ----
        desc = (
            f"slide {slide_index + 1}: unsupported shape type "
            f"{shape_type!r} ('{shape.name}')"
        )
        warnings.append(desc)
        blocks.append(UnsupportedBlock(description=desc))

    # ------------------------------------------------------------------
    # Text frame extraction
    # ------------------------------------------------------------------

    def _extract_text_frame(self, shape, blocks: list[Block]) -> None:
        """Emit :class:`ParagraphBlock` entries for each non-empty paragraph."""
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                blocks.append(ParagraphBlock(text=text))

    # ------------------------------------------------------------------
    # Picture / raster image extraction
    # ------------------------------------------------------------------

    def _extract_picture(
        self,
        shape,
        embedded: list[EmbeddedAsset],
        warnings: list[str],
    ) -> ImageBlock | None:
        """Extract a raster image from a picture shape."""
        try:
            image = shape.image
            content_type: str = image.content_type or ""
            ext = _content_type_to_ext(content_type)

            # Check if this is actually a vector format stored as a picture
            src_fmt = _content_type_to_vector_format(content_type)
            if src_fmt is not None:
                return self._convert_vector_asset(
                    image.blob, src_fmt, embedded, warnings
                )

            original_name = Path(image.filename).name if image.filename else None
            alt = original_name or shape.name or "image"
            asset = EmbeddedAsset(
                data=image.blob,
                extension=ext,
                original_name=original_name,
                alt_text=alt,
                source_vector_format=None,
            )
            idx = len(embedded)
            embedded.append(asset)
            return ImageBlock(asset_index=idx, alt=alt)
        except Exception as exc:  # noqa: BLE001
            warnings.append(f"failed to extract picture '{shape.name}': {exc}")
            return None

    # ------------------------------------------------------------------
    # OLE object extraction (may contain EMF/WMF)
    # ------------------------------------------------------------------

    def _extract_ole_object(
        self,
        shape,
        embedded: list[EmbeddedAsset],
        warnings: list[str],
    ) -> ImageBlock | None:
        """Try to extract a vector graphic from an OLE object shape."""
        try:
            # python-pptx exposes the image of an OLE object via shape.image
            # when the OLE object has a picture representation
            image = shape.image
            content_type: str = image.content_type or ""
            src_fmt = _content_type_to_vector_format(content_type)
            if src_fmt is not None:
                return self._convert_vector_asset(
                    image.blob, src_fmt, embedded, warnings
                )
            # Raster fallback for OLE with non-vector image
            ext = _content_type_to_ext(content_type)
            original_name = Path(image.filename).name if image.filename else None
            alt = original_name or shape.name or "ole object"
            asset = EmbeddedAsset(
                data=image.blob,
                extension=ext,
                original_name=original_name,
                alt_text=alt,
                source_vector_format=None,
            )
            idx = len(embedded)
            embedded.append(asset)
            return ImageBlock(asset_index=idx, alt=alt)
        except Exception as exc:  # noqa: BLE001
            logger.debug("OLE object extraction failed for '%s': %s", shape.name, exc)
            return None

    # ------------------------------------------------------------------
    # Vector asset conversion
    # ------------------------------------------------------------------

    def _convert_vector_asset(
        self,
        data: bytes,
        src_fmt: str,
        embedded: list[EmbeddedAsset],
        warnings: list[str],
    ) -> ImageBlock | None:
        """Run *data* through :class:`VectorConverter` and emit an asset."""
        try:
            out_bytes, ext = self._vector.convert(data, src_fmt)  # type: ignore[arg-type]
            asset = EmbeddedAsset(
                data=out_bytes,
                extension=ext,
                original_name=None,
                alt_text=f"{src_fmt} drawing",
                source_vector_format=src_fmt,
            )
            idx = len(embedded)
            embedded.append(asset)
            return ImageBlock(asset_index=idx, alt=asset.alt_text)
        except VectorConversionError as exc:
            warnings.append(f"vector conversion failed ({src_fmt}): {exc}")
            return None
