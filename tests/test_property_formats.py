"""Property-based tests for supported format conversion.

# Feature: document-to-markdown, Property 1: Supported format produces output file
# Feature: document-to-markdown, Property 3: Output is valid UTF-8
"""
from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
from hypothesis import given, settings
from hypothesis import strategies as st

from document2markdown.api import Converter
from document2markdown.renderer_base import MarkdownRenderer

# ---------------------------------------------------------------------------
# Minimal valid document builders
# ---------------------------------------------------------------------------

def _make_txt(tmp_dir: Path, name: str = "sample") -> Path:
    p = tmp_dir / f"{name}.txt"
    p.write_text("Hello, world! This is a plain text document.", encoding="utf-8")
    return p


def _make_html(tmp_dir: Path, name: str = "sample") -> Path:
    p = tmp_dir / f"{name}.html"
    p.write_text(
        "<html><body><h1>Title</h1><p>A paragraph.</p></body></html>",
        encoding="utf-8",
    )
    return p


def _make_docx(tmp_dir: Path, name: str = "sample") -> Path:
    from docx import Document as DocxDocument  # type: ignore

    doc = DocxDocument()
    doc.add_heading("Sample Heading", level=1)
    doc.add_paragraph("A sample paragraph.")
    p = tmp_dir / f"{name}.docx"
    doc.save(str(p))
    return p


def _make_pptx(tmp_dir: Path, name: str = "sample") -> Path:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches  # type: ignore

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    if title:
        title.text = "Sample Slide"
    p = tmp_dir / f"{name}.pptx"
    prs.save(str(p))
    return p


def _make_pdf(tmp_dir: Path, name: str = "sample") -> Path:
    """Create a minimal valid PDF using PyMuPDF (fitz)."""
    import fitz  # type: ignore  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Hello PDF world!")
    p = tmp_dir / f"{name}.pdf"
    doc.save(str(p))
    doc.close()
    return p


# Map format name → builder function
_FORMAT_BUILDERS = {
    "txt":  _make_txt,
    "html": _make_html,
    "docx": _make_docx,
    "pptx": _make_pptx,
    "pdf":  _make_pdf,
}

_FORMAT_NAMES = sorted(_FORMAT_BUILDERS.keys())


# ---------------------------------------------------------------------------
# Property 1: Supported format produces output file
# ---------------------------------------------------------------------------

# Feature: document-to-markdown, Property 1: Supported format produces output file
@given(fmt=st.sampled_from(_FORMAT_NAMES))
@settings(max_examples=100)
def test_property1_supported_format_produces_output_file(fmt: str) -> None:
    """Converting a minimal valid document of each supported format produces a non-empty .md file."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp = Path(tmp_dir)
        src_path = _FORMAT_BUILDERS[fmt](tmp)
        out_dir = tmp / "output"
        out_dir.mkdir()

        converter = Converter(output_dir=out_dir)
        doc = converter.convert(src_path)
        doc.save(out_dir)

        expected_md = out_dir / f"{src_path.stem}.md"
        assert expected_md.exists(), f"Output .md file not found: {expected_md}"
        content = expected_md.read_text(encoding="utf-8")
        assert content.strip(), f"Output .md file is empty for format {fmt!r}"


# ---------------------------------------------------------------------------
# Property 3: Output is valid UTF-8
# ---------------------------------------------------------------------------

# Feature: document-to-markdown, Property 3: Output is valid UTF-8
@given(fmt=st.sampled_from(_FORMAT_NAMES))
@settings(max_examples=100)
def test_property3_output_is_valid_utf8(fmt: str) -> None:
    """The rendered Markdown string for any supported format is valid UTF-8."""
    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp = Path(tmp_dir)
        src_path = _FORMAT_BUILDERS[fmt](tmp)

        converter = Converter()
        doc = converter.convert(src_path)
        markdown = doc.to_markdown()

        # Must encode and decode as UTF-8 without error
        try:
            encoded = markdown.encode("utf-8")
            decoded = encoded.decode("utf-8")
        except (UnicodeEncodeError, UnicodeDecodeError) as exc:
            pytest.fail(
                f"Output for format {fmt!r} is not valid UTF-8: {exc}\n"
                f"Markdown snippet: {markdown[:200]!r}"
            )

        assert decoded == markdown, "UTF-8 round-trip changed the string"
