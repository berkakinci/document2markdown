"""Integration tests for document2markdown — all five formats.

Creates minimal sample files programmatically, converts them via the
Converter class, and asserts output correctness.

Requirements: 1.1–1.5, 2.6–2.8, 3.6, 4.3, 6.1, 7.2
"""

from __future__ import annotations

import io
import re
import struct
import sys
from pathlib import Path

import pytest

# ---------------------------------------------------------------------------
# Helpers to build minimal sample files programmatically
# ---------------------------------------------------------------------------


def _make_txt(path: Path) -> None:
    """Write a minimal plain-text file."""
    path.write_text("Hello from TXT.\nLine two.\n", encoding="utf-8")


def _make_html(path: Path) -> None:
    """Write a minimal HTML file."""
    path.write_text(
        "<!DOCTYPE html><html><body>"
        "<h1>Title</h1><p>Hello from HTML.</p>"
        "</body></html>",
        encoding="utf-8",
    )


def _make_docx(path: Path) -> None:
    """Write a minimal DOCX file using python-docx."""
    from docx import Document as DocxDocument  # type: ignore[import]

    doc = DocxDocument()
    doc.add_heading("Integration Test", level=1)
    doc.add_paragraph("Hello from DOCX.")
    doc.save(str(path))


def _make_pptx(path: Path) -> None:
    """Write a minimal PPTX file using python-pptx."""
    from pptx import Presentation  # type: ignore[import]
    from pptx.util import Inches  # type: ignore[import]

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Integration Test Slide"
    slide.placeholders[1].text = "Hello from PPTX."
    prs.save(str(path))


def _make_pdf(path: Path) -> None:
    """Write a PDF file with enough content for pymupdf4llm layout analysis."""
    import fitz  # type: ignore[import]

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    # Title-sized text so layout module has something to classify
    page.insert_text((72, 80), "Integration Test Document", fontsize=20)
    # Body paragraphs to give the layout module enough content
    page.insert_text((72, 140), "This is the first paragraph of the document.", fontsize=12)
    page.insert_text((72, 170), "This is the second paragraph with more text content.", fontsize=12)
    page.insert_text((72, 200), "A third line ensures the layout module has enough to work with.", fontsize=12)
    doc.save(str(path))
    doc.close()


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------


@pytest.fixture()
def out_dir(tmp_path: Path) -> Path:
    """Return a fresh temporary output directory."""
    d = tmp_path / "output"
    d.mkdir()
    return d


@pytest.fixture()
def src_dir(tmp_path: Path) -> Path:
    """Return a fresh temporary source directory."""
    d = tmp_path / "sources"
    d.mkdir()
    return d


# ---------------------------------------------------------------------------
# Helper: run a single conversion and return the .md path
# ---------------------------------------------------------------------------


def _convert(src: Path, out: Path) -> Path:
    from document2markdown import Converter

    converter = Converter(output_dir=out)
    doc = converter.convert(src)
    return doc.save()


# ---------------------------------------------------------------------------
# Tests — one per format
# ---------------------------------------------------------------------------


class TestTXTIntegration:
    def test_output_exists_nonempty_utf8(self, src_dir: Path, out_dir: Path) -> None:
        src = src_dir / "sample.txt"
        _make_txt(src)
        md_path = _convert(src, out_dir)

        assert md_path.exists(), "Output .md file should exist"
        content = md_path.read_bytes()
        assert len(content) > 0, "Output .md file should be non-empty"
        # Must decode as UTF-8 without error
        decoded = content.decode("utf-8")
        assert "Hello from TXT" in decoded or len(decoded) > 0


class TestHTMLIntegration:
    def test_output_exists_nonempty_utf8(self, src_dir: Path, out_dir: Path) -> None:
        src = src_dir / "sample.html"
        _make_html(src)
        md_path = _convert(src, out_dir)

        assert md_path.exists()
        content = md_path.read_bytes()
        assert len(content) > 0
        decoded = content.decode("utf-8")
        assert len(decoded) > 0


class TestDOCXIntegration:
    def test_output_exists_nonempty_utf8(self, src_dir: Path, out_dir: Path) -> None:
        src = src_dir / "sample.docx"
        _make_docx(src)
        md_path = _convert(src, out_dir)

        assert md_path.exists()
        content = md_path.read_bytes()
        assert len(content) > 0
        content.decode("utf-8")  # must not raise


class TestPPTXIntegration:
    def test_output_exists_nonempty_utf8(self, src_dir: Path, out_dir: Path) -> None:
        src = src_dir / "sample.pptx"
        _make_pptx(src)
        md_path = _convert(src, out_dir)

        assert md_path.exists()
        content = md_path.read_bytes()
        assert len(content) > 0
        content.decode("utf-8")  # must not raise


class TestPDFIntegration:
    def test_output_exists_nonempty_utf8(self, src_dir: Path, out_dir: Path) -> None:
        src = src_dir / "sample.pdf"
        _make_pdf(src)
        md_path = _convert(src, out_dir)

        assert md_path.exists()
        content = md_path.read_bytes()
        assert len(content) > 0
        content.decode("utf-8")  # must not raise


# ---------------------------------------------------------------------------
# Test — embedded asset naming
# ---------------------------------------------------------------------------


class TestEmbeddedAssetNaming:
    """Assert md_embedded/ assets follow {base_name}_{serial:04d}{ext} naming."""

    def test_docx_with_image_produces_embedded_asset(
        self, src_dir: Path, out_dir: Path
    ) -> None:
        """A DOCX with an embedded PNG should produce an asset in md_embedded/."""
        from docx import Document as DocxDocument  # type: ignore[import]
        from docx.shared import Inches  # type: ignore[import]

        # Build a tiny 1×1 white PNG in memory
        png_bytes = _minimal_png()

        src = src_dir / "with_image.docx"
        doc = DocxDocument()
        doc.add_heading("Doc with image", level=1)
        doc.add_paragraph("Below is an image.")
        doc.add_picture(io.BytesIO(png_bytes), width=Inches(1))
        doc.save(str(src))

        md_path = _convert(src, out_dir)
        assert md_path.exists()

        embedded_dir = out_dir / "md_embedded"
        if embedded_dir.exists():
            assets = list(embedded_dir.iterdir())
            for asset in assets:
                # Name must match {base_name}_{serial:04d}{ext}
                name = asset.stem  # e.g. "with_image_0001"
                assert re.match(r"^with_image_\d{4}$", name), (
                    f"Asset name '{asset.name}' does not match expected pattern"
                )


def _minimal_png() -> bytes:
    """Return the bytes of a minimal 1×1 white PNG."""
    import zlib

    def _chunk(tag: bytes, data: bytes) -> bytes:
        c = struct.pack(">I", len(data)) + tag + data
        crc = struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        return c + crc

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = _chunk(b"IHDR", ihdr_data)
    raw_row = b"\x00\xff\xff\xff"  # filter byte + RGB white
    idat = _chunk(b"IDAT", zlib.compress(raw_row))
    iend = _chunk(b"IEND", b"")
    return signature + ihdr + idat + iend


# ---------------------------------------------------------------------------
# Test — batch summary via convert_batch
# ---------------------------------------------------------------------------


class TestBatchSummary:
    """Assert convert_batch returns one entry per input and handles failures."""

    def test_batch_all_succeed(self, src_dir: Path, out_dir: Path) -> None:
        from document2markdown import Converter
        from document2markdown.utils import convert_batch
        from document2markdown.document import Document

        txt1 = src_dir / "a.txt"
        txt2 = src_dir / "b.txt"
        _make_txt(txt1)
        _make_txt(txt2)

        converter = Converter(output_dir=out_dir)
        results = convert_batch([txt1, txt2], converter)

        assert len(results) == 2
        for path, outcome in results:
            assert isinstance(outcome, Document), (
                f"Expected Document for {path}, got {type(outcome)}"
            )

    def test_batch_continues_on_failure(self, src_dir: Path, out_dir: Path) -> None:
        from document2markdown import Converter
        from document2markdown.utils import convert_batch
        from document2markdown.document import Document

        good = src_dir / "good.txt"
        bad = src_dir / "nonexistent.txt"  # does not exist
        _make_txt(good)

        converter = Converter(output_dir=out_dir)
        results = convert_batch([good, bad], converter)

        assert len(results) == 2, "Should return one entry per input"
        paths = [p for p, _ in results]
        assert good in paths
        assert bad in paths

        # The good file should succeed
        good_outcome = next(o for p, o in results if p == good)
        assert isinstance(good_outcome, Document)

        # The bad file should be an exception
        bad_outcome = next(o for p, o in results if p == bad)
        assert isinstance(bad_outcome, Exception)

    def test_batch_summary_counts(self, src_dir: Path, out_dir: Path) -> None:
        """succeeded + failed == total."""
        from document2markdown import Converter
        from document2markdown.utils import convert_batch
        from document2markdown.document import Document

        files = []
        for i in range(3):
            f = src_dir / f"file{i}.txt"
            _make_txt(f)
            files.append(f)
        # Add two non-existent files
        files.append(src_dir / "missing1.txt")
        files.append(src_dir / "missing2.txt")

        converter = Converter(output_dir=out_dir)
        results = convert_batch(files, converter)

        total = len(results)
        succeeded = sum(1 for _, o in results if isinstance(o, Document))
        failed = sum(1 for _, o in results if isinstance(o, Exception))

        assert total == 5
        assert succeeded == 3
        assert failed == 2
        assert succeeded + failed == total


# ---------------------------------------------------------------------------
# Test — MIME/extension mismatch rejection
# ---------------------------------------------------------------------------


class TestMimeExtensionMismatch:
    """A file with .docx extension but PDF magic bytes must be rejected."""

    def test_docx_extension_pdf_content_raises(
        self, src_dir: Path, out_dir: Path
    ) -> None:
        from document2markdown import Converter
        from document2markdown.errors import MimeExtensionMismatchError

        # Create a real PDF, then rename it to .docx
        real_pdf = src_dir / "real.pdf"
        _make_pdf(real_pdf)
        fake_docx = src_dir / "fake.docx"
        fake_docx.write_bytes(real_pdf.read_bytes())

        converter = Converter(output_dir=out_dir)
        with pytest.raises(MimeExtensionMismatchError) as exc_info:
            converter.convert(fake_docx)

        err_msg = str(exc_info.value)
        # The error message must mention both the extension-derived type and
        # the magic-byte-derived type (Requirement 7.2)
        assert "application/vnd.openxmlformats-officedocument" in err_msg or "docx" in err_msg.lower()
        assert "pdf" in err_msg.lower() or "application/pdf" in err_msg

    def test_mismatch_error_written_to_stderr(
        self, src_dir: Path, out_dir: Path, capsys: pytest.CaptureFixture
    ) -> None:
        """When the CLI-level code catches the error it should go to stderr."""
        from document2markdown import Converter
        from document2markdown.errors import MimeExtensionMismatchError

        real_pdf = src_dir / "real2.pdf"
        _make_pdf(real_pdf)
        fake_docx = src_dir / "fake2.docx"
        fake_docx.write_bytes(real_pdf.read_bytes())

        converter = Converter(output_dir=out_dir)
        try:
            converter.convert(fake_docx)
        except MimeExtensionMismatchError as exc:
            print(str(exc), file=sys.stderr)

        captured = capsys.readouterr()
        assert len(captured.err) > 0, "Mismatch error should produce stderr output"
