"""Unit tests for PPTXConverter — Task 6.2."""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from document2markdown.converter_pptx import PPTXConverter
from document2markdown.document_model import HeadingBlock, ParagraphBlock


def _make_pptx(tmp_path: Path, builder) -> Path:
    prs = Presentation()
    builder(prs)
    p = tmp_path / "test.pptx"
    prs.save(str(p))
    return p


class TestSlideTitle:
    def test_title_becomes_h2(self, tmp_path):
        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "My Slide Title"
            slide.placeholders[1].text = "body text"

        p = _make_pptx(tmp_path, build)
        result = PPTXConverter().convert(p)
        headings = [b for b in result.blocks if isinstance(b, HeadingBlock)]
        assert headings, "Expected at least one HeadingBlock"
        assert headings[0].level == 2
        assert headings[0].text == "My Slide Title"


class TestBodyText:
    def test_body_text_becomes_paragraph(self, tmp_path):
        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Title"
            slide.placeholders[1].text = "Body content here."

        p = _make_pptx(tmp_path, build)
        result = PPTXConverter().convert(p)
        paras = [b for b in result.blocks if isinstance(b, ParagraphBlock)]
        assert any("Body content here." in b.text for b in paras)


class TestMultipleSlides:
    def test_slides_processed_in_order(self, tmp_path):
        def build(prs):
            for title in ("Slide One", "Slide Two", "Slide Three"):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = title
                slide.placeholders[1].text = f"Content of {title}"

        p = _make_pptx(tmp_path, build)
        result = PPTXConverter().convert(p)
        headings = [b for b in result.blocks if isinstance(b, HeadingBlock)]
        titles = [h.text for h in headings]
        assert titles == ["Slide One", "Slide Two", "Slide Three"]
        assert all(h.level == 2 for h in headings)


class TestEmptyParagraphsFiltered:
    def test_blank_paragraphs_not_emitted(self, tmp_path):
        """Empty paragraphs inside a text frame should not produce blocks."""
        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
            tf = txBox.text_frame
            tf.text = "First paragraph"
            tf.add_paragraph()  # blank
            p = tf.add_paragraph()
            p.text = "Third paragraph"

        p = _make_pptx(tmp_path, build)
        result = PPTXConverter().convert(p)
        paras = [b for b in result.blocks if isinstance(b, ParagraphBlock)]
        texts = [b.text for b in paras]
        assert "First paragraph" in texts
        assert "Third paragraph" in texts
        assert "" not in texts
        assert any(t == "" for t in texts) is False


class TestGroupShape:
    def test_group_shape_children_extracted(self, tmp_path):
        """Text inside a group shape should be extracted via recursion."""
        from pptx.util import Inches, Pt

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            # Add two text boxes and group them
            tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
            tb1.text_frame.text = "Group child one"
            tb2 = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
            tb2.text_frame.text = "Group child two"
            # python-pptx doesn't expose group creation directly, so we test
            # the recursion path by mocking a group shape
            pass  # see mock-based test below

        # Use a mock to exercise the GROUP branch directly
        from unittest.mock import MagicMock, patch
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        converter = PPTXConverter()
        blocks = []
        embedded = []
        warnings = []

        child1 = MagicMock()
        child1.has_text_frame = True
        child1.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        child1.text_frame.paragraphs = [MagicMock(text="  Child A  ")]
        child1.text_frame.paragraphs[0].text = "Child A"

        child2 = MagicMock()
        child2.has_text_frame = True
        child2.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
        child2.text_frame.paragraphs = [MagicMock(text="Child B")]
        child2.text_frame.paragraphs[0].text = "Child B"

        group = MagicMock()
        group.has_text_frame = False
        group.shape_type = MSO_SHAPE_TYPE.GROUP
        group.shapes = [child1, child2]

        converter._convert_shape(group, 0, blocks, embedded, warnings)

        texts = [b.text for b in blocks if isinstance(b, ParagraphBlock)]
        assert "Child A" in texts
        assert "Child B" in texts


class TestAutoShapeSilentSkip:
    def test_autoshape_produces_no_block(self, tmp_path):
        """AUTO_SHAPE, FREEFORM, and LINE shapes should be silently skipped."""
        from unittest.mock import MagicMock
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        converter = PPTXConverter()

        for shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.LINE,
        ):
            blocks = []
            shape = MagicMock()
            shape.has_text_frame = False
            shape.shape_type = shape_type
            converter._convert_shape(shape, 0, blocks, [], [])
            assert blocks == [], f"Expected no blocks for shape_type={shape_type}"


class TestUnsupportedShapeType:
    def test_unknown_shape_emits_unsupported_block(self, tmp_path):
        """An unrecognised shape type should produce an UnsupportedBlock and a warning."""
        from unittest.mock import MagicMock
        from document2markdown.document_model import UnsupportedBlock

        converter = PPTXConverter()
        blocks = []
        warnings = []

        shape = MagicMock()
        shape.has_text_frame = False
        shape.shape_type = 9999  # not a real MSO_SHAPE_TYPE
        shape.name = "WeirdShape"

        converter._convert_shape(shape, 0, blocks, [], warnings)

        assert len(blocks) == 1
        assert isinstance(blocks[0], UnsupportedBlock)
        assert len(warnings) == 1


class TestPictureExtraction:
    def test_picture_shape_produces_image_block(self, tmp_path):
        """A PICTURE shape should produce an ImageBlock and an EmbeddedAsset."""
        from unittest.mock import MagicMock, PropertyMock
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from document2markdown.document_model import ImageBlock

        converter = PPTXConverter()
        blocks = []
        embedded = []
        warnings = []

        # Minimal 1x1 PNG bytes
        png_bytes = (
            b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01'
            b'\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00'
            b'\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18'
            b'\xd8N\x00\x00\x00\x00IEND\xaeB`\x82'
        )

        image = MagicMock()
        image.content_type = "image/png"
        image.blob = png_bytes
        image.filename = "photo.png"

        shape = MagicMock()
        shape.has_text_frame = False
        shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        shape.image = image
        shape.name = "Picture 1"

        converter._convert_shape(shape, 0, blocks, embedded, warnings)

        assert len(blocks) == 1
        assert isinstance(blocks[0], ImageBlock)
        assert len(embedded) == 1
        assert embedded[0].extension == ".png"
        assert warnings == []

    def test_picture_extraction_failure_adds_warning(self, tmp_path):
        """If image extraction raises, a warning is added and no block emitted."""
        from unittest.mock import MagicMock, PropertyMock
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        converter = PPTXConverter()
        blocks = []
        warnings = []

        shape = MagicMock()
        shape.has_text_frame = False
        shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        shape.name = "BadPic"
        type(shape).image = PropertyMock(side_effect=Exception("no image data"))

        converter._convert_shape(shape, 0, blocks, [], warnings)

        assert blocks == []
        assert len(warnings) == 1
        assert "BadPic" in warnings[0]


class TestOLEObjectExtraction:
    def test_ole_failure_emits_unsupported_block(self, tmp_path):
        """OLE object that can't be extracted should emit UnsupportedBlock + warning."""
        from unittest.mock import MagicMock, PropertyMock
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from document2markdown.document_model import UnsupportedBlock

        converter = PPTXConverter()
        blocks = []
        warnings = []

        shape = MagicMock()
        shape.has_text_frame = False
        shape.shape_type = MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
        shape.name = "OLEObj"
        type(shape).image = PropertyMock(side_effect=Exception("no image"))

        converter._convert_shape(shape, 0, blocks, [], warnings)

        assert len(blocks) == 1
        assert isinstance(blocks[0], UnsupportedBlock)
        assert len(warnings) == 1
